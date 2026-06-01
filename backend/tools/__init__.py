import json, os, subprocess, requests, time, asyncio, threading
from datetime import datetime
from config import AION_DIR, MEMORY_FILE, MEMORY_DIR, ACTIVITY_FILE, UPLOADS_DIR, SESSIONS_DIR, CRM_DIR, LEADS_FILE

PERF_START = {}
MEMORY_LOCK = threading.Lock()

# Memory recall cache (30s TTL)
_recall_cache = {}
RECALL_CACHE_TTL = 30

def _get_memory_file():
    try:
        from kb import _get_current_project
        project = _get_current_project()
        if project and project != "default":
            pfile = os.path.join(AION_DIR, "MEMORY", project, "memory.json")
            if os.path.exists(pfile):
                return pfile
            legacy = os.path.join(AION_DIR, "MEMORY", "memory.json")
            if os.path.exists(legacy):
                os.makedirs(os.path.dirname(pfile), exist_ok=True)
                import shutil
                shutil.copy2(legacy, pfile)
                return pfile
            return pfile
    except:
        pass
    return MEMORY_FILE

def log_activity(agent_id, tool, args, result, success=True):
    try:
        os.makedirs(os.path.dirname(ACTIVITY_FILE), exist_ok=True)
        entry = {
            "ts": datetime.now().isoformat(),
            "agent": agent_id,
            "tool": tool,
            "args": str(args)[:200],
            "result": str(result)[:300],
            "success": success,
        }
        with open(ACTIVITY_FILE, "a") as f:
            f.write(json.dumps(entry, ensure_ascii=False) + "\n")
    except:
        pass

def read_activity(limit=100):
    try:
        if not os.path.exists(ACTIVITY_FILE):
            return []
        entries = []
        with open(ACTIVITY_FILE) as f:
            for line in f:
                line = line.strip()
                if line:
                    try:
                        entries.append(json.loads(line))
                    except:
                        pass
        return entries[-limit:]
    except:
        return []

def store_collab_memory(agent_id, task, result):
    mem = load_memory()
    if "collaborations" not in mem:
        mem["collaborations"] = []
    mem["collaborations"].append({
        "agent": agent_id,
        "task": task[:300],
        "result": result[:500],
        "timestamp": datetime.now().isoformat(),
    })
    if len(mem["collaborations"]) > 100:
        mem["collaborations"] = mem["collaborations"][-100:]
    save_memory(mem)

def load_memory():
    mfile = _get_memory_file()
    with MEMORY_LOCK:
        try:
            os.makedirs(os.path.dirname(mfile), exist_ok=True)
            if os.path.exists(mfile):
                with open(mfile) as f:
                    return json.load(f)
        except:
            pass
        return {}

def save_memory(data):
    mfile = _get_memory_file()
    with MEMORY_LOCK:
        try:
            os.makedirs(os.path.dirname(mfile), exist_ok=True)
            with open(mfile, "w") as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
            return "OK"
        except Exception as e:
            return f"σφάλμα: {e}"

def _get_facts(mem):
    return mem.get("facts", {})

def _ensure_facts(mem):
    if "facts" not in mem:
        mem["facts"] = {}
    return mem["facts"]

def get_tool_definitions_for_agent(agent_id):
    from agents import AGENTS
    agent_tools = []
    for a in AGENTS:
        if a["id"] == agent_id:
            agent_tools = a.get("tools", [])
            break
    if not agent_tools:
        return TOOL_DEFINITIONS
    return [td for td in TOOL_DEFINITIONS if td["function"]["name"] in agent_tools]

TOOL_DEFINITIONS = [
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Διάβασε περιεχόμενο αρχείου. Path must be absolute or under ~/AION.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Απόλυτο path ή σχετικό με ~/AION"}
                },
                "required": ["path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Γράψε περιεχόμενο (TEXT/HTML/markdown) σε αρχείο. Δημιουργεί το αρχείο αν δεν υπάρχει. ΜΗΝ το χρησιμοποιείς για Word/Excel/PowerPoint — υπάρχουν generate_docx, generate_xlsx, generate_pptx γι' αυτά.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Απόλυτο path ή σχετικό με ~/AION"},
                    "content": {"type": "string", "description": "Περιεχόμενο προς εγγραφή"}
                },
                "required": ["path", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_dir",
            "description": "Λίστα αρχείων σε directory.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Απόλυτο path ή σχετικό με ~/AION"}
                },
                "required": ["path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": "Εκτέλεσε command στο terminal.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "Command προς εκτέλεση"},
                    "timeout": {"type": "number", "description": "Timeout σε δευτερόλεπτα"}
                },
                "required": ["command"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "Αναζήτηση στο web.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Search query"}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "web_fetch",
            "description": "Διάβασε περιεχόμενο από URL.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "URL προς ανάγνωση"},
                    "format": {"type": "string", "enum": ["markdown", "text", "html"], "description": "Μορφή εξόδου"}
                },
                "required": ["url"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "remember",
            "description": "Αποθήκευσε μια πληροφορία στη μνήμη.",
            "parameters": {
                "type": "object",
                "properties": {
                    "key": {"type": "string", "description": "Όνομα της πληροφορίας"},
                    "value": {"type": "string", "description": "Τιμή"}
                },
                "required": ["key", "value"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "recall",
            "description": "Ανάκτησε πληροφορία από τη μνήμη.",
            "parameters": {
                "type": "object",
                "properties": {
                    "key": {"type": "string", "description": "Όνομα προς αναζήτηση"}
                },
                "required": ["key"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_memories",
            "description": "Δες όλες τις αποθηκευμένες πληροφορίες.",
            "parameters": {
                "type": "object",
                "properties": {}
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_leads",
            "description": "Διάβασε leads από το CRM database.",
            "parameters": {
                "type": "object",
                "properties": {
                    "search": {"type": "string", "description": "Προαιρετική αναζήτηση"},
                    "limit": {"type": "number", "description": "Αριθμός leads (default 10)"}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "save_lead",
            "description": "Αποθήκευσε ένα νέο lead στο CRM database. ΧΡΗΣΙΜΟΠΟΙΗΣΕ αυτό το tool ΑΦΟΥ κάνεις web_search για πραγματικά leads. ΠΡΕΠΕΙ να δώσεις πραγματικά στοιχεία που βρήκες από web search, όχι επινοημένα.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Όνομα επιχείρησης (υποχρεωτικό)"},
                    "website": {"type": "string", "description": "Ιστοσελίδα"},
                    "industry": {"type": "string", "description": "Κλάδος"},
                    "location": {"type": "string", "description": "Τοποθεσία"},
                    "contact": {"type": "string", "description": "Στοιχεία επικοινωνίας"},
                    "serviceNeeded": {"type": "string", "description": "Υπηρεσία που χρειάζεται"},
                    "onlinePresence": {"type": "string", "description": "none / basic / moderate / strong"},
                    "source_url": {"type": "string", "description": "URL όπου βρήκες το lead"},
                    "notes": {"type": "string", "description": "Σημειώσεις"}
                },
                "required": ["name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_time",
            "description": "Τρέχουσα ώρα και ημερομηνία.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_to_agent",
            "description": "ΣΤΕΙΛΕ μήνυμα σε άλλο agent. Η επικοινωνία μεταξύ agents γίνεται σε COMPACT format (δομημένο, όχι φυσική γλώσσα) για ταχύτητα. Χρησιμοποίησε format:'full' μόνο για τελικά μηνύματα προς τον χρήστη.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent_id": {
                        "type": "string",
                        "enum": ["dev", "leadfinder", "memory", "sales", "marketing", "content", "support", "analytics", "security", "finance", "imggen", "seo", "offers", "pm", "consultant", "docsagent", "ceo"],
                        "description": "Ποιος agent θα λάβει το μήνυμα"
                    },
                    "message": {"type": "string", "description": "Το μήνυμα — σε COMPACT format για agent-to-agent, φυσική γλώσσα μόνο αν format:'full'"},
                    "context": {"type": "string", "description": "Πρόσθετες πληροφορίες context"},
                    "format": {
                        "type": "string",
                        "enum": ["compact", "full"],
                        "description": "compact: σύντομο/δομημένο (default, για agent-to-agent). full: φυσική γλώσσα (για τελικό output στον χρήστη)"
                    }
                },
                "required": ["agent_id", "message"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_file_to_agent",
            "description": "ΣΤΕΙΛΕ ένα αρχείο σε άλλο agent. Το αρχείο θα αντιγραφεί στο φάκελο του παραλήπτη και θα είναι διαθέσιμο προς ανάγνωση. Χρησιμοποίησέ το για να μοιραστείς αποτελέσματα web search, reports, ή άλλα αρχεία.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent_id": {
                        "type": "string",
                        "enum": ["ceo", "dev", "leadfinder", "memory", "sales", "marketing", "content", "support", "analytics", "security", "finance", "imggen", "seo", "offers", "pm", "consultant", "docsagent"],
                        "description": "Σε ποιον agent να σταλεί το αρχείο"
                    },
                    "file_path": {"type": "string", "description": "Απόλυτο path του αρχείου προς αποστολή"},
                    "rename": {"type": "string", "description": "Προαιρετικό: νέο όνομα για το αρχείο στον παραλήπτη"}
                },
                "required": ["agent_id", "file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_agents",
            "description": "Δες όλους τους διαθέσιμους agents στο σύστημα, τις δυνατότητές τους και τα εργαλεία τους.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "delegate_to_agent",
            "description": "ΑΝΑΘΕΣΕ εργασία σε άλλο agent σε COMPACT format. Μόνο το τελικό output προς τον χρήστη γράφεται σε φυσική γλώσσα.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent_id": {
                        "type": "string",
                        "enum": ["dev", "leadfinder", "memory", "sales", "marketing", "content", "support", "analytics", "security", "finance", "imggen", "seo", "offers", "pm", "consultant", "docsagent"],
                        "description": "Ποιος agent θα εκτελέσει την εργασία"
                    },
                    "task": {"type": "string", "description": "COMPACT task description (όχι φυσική γλώσσα — π.χ. AUDIT src/auth.py FOCUS: sqli,xss)"},
                    "context": {"type": "string", "description": "Πρόσθετες πληροφορίες (compact)"},
                    "format": {
                        "type": "string",
                        "enum": ["compact", "full"],
                        "description": "compact: δομημένο/σύντομο (default). full: φυσική γλώσσα (μόνο όταν το αποτέλεσμα πάει σε χρήστη)"
                    }
                },
                "required": ["agent_id", "task"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "parallel_delegate",
            "description": "ΑΝΑΘΕΣΕ εργασίες σε ΠΟΛΛΟΥΣ agents ταυτόχρονα σε COMPACT format. Οι agents τρέχουν παράλληλα και επιστρέφουν αποτελέσματα.",
            "parameters": {
                "type": "object",
                "properties": {
                    "delegations": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "agent_id": {
                                    "type": "string",
                                    "enum": ["dev", "leadfinder", "memory", "sales", "marketing", "content", "support", "analytics", "security", "finance", "imggen", "seo", "offers", "pm", "consultant", "docsagent"],
                                    "description": "Ποιος agent θα εκτελέσει την εργασία"
                                },
                                "task": {"type": "string", "description": "COMPACT task (όχι φυσική γλώσσα)"},
                                "context": {"type": "string", "description": "Πρόσθετες πληροφορίες (compact, προαιρετικό)"},
                                "format": {
                                    "type": "string",
                                    "enum": ["compact", "full"],
                                    "description": "compact: δομημένο/σύντομο (default). full: φυσική γλώσσα"
                                }
                            },
                            "required": ["agent_id", "task"]
                        },
                        "description": "Λίστα από εργασίες προς εκτέλεση παράλληλα"
                    },
                    "synthesize": {"type": "boolean", "description": "Αν θέλεις να συνθέσεις τα αποτελέσματα σε ενιαία απάντηση (true) ή να τα επιστρέψεις ξεχωριστά (false, default)"}
                },
                "required": ["delegations"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "query_kb",
            "description": "Αναζήτησε στο Knowledge Base του project. Επιστρέφει σχετικά αποσπάσματα από αρχεία και σημειώσεις που έχουν αποθηκευτεί στο project ή στο global knowledge base.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Ερώτημα αναζήτησης (π.χ. 'brand guidelines', 'API endpoints', 'προσφορές')"},
                    "project": {"type": "string", "description": "Project name (προαιρετικό, default το τρέχον project)"}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_agent_history",
            "description": "Διάβασε το ιστορικό συνομιλιών ενός agent. Επιστρέφει τα τελευταία μηνύματα για να δεις τι έχει κάνει ο agent.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent": {"type": "string", "description": "ID του agent (π.χ. dev, sales, memory)"}
                },
                "required": ["agent"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "generate_docx",
            "description": "Δημιούργησε ΕΠΑΓΓΕΛΜΑΤΙΚΟ Word έγγραφο (.docx) πολλών σελίδων με εκτενή ανάλυση, πολλαπλές ενότητες, πίνακες, λίστες. ΑΠΑΙΤΟΥΝΤΑΙ πολλές λεπτομέρειες — το έγγραφο πρέπει να είναι πλήρες, αναλυτικό, επαγγελματικό. Ιδανικό για proposals, reports, συμβόλαια, προσφορές. ΧΡΗΣΙΜΟΠΟΙΗΣΕ το για any professional έγγραφο — μη γράφεις απλό κείμενο. ΠΡΟΣΟΧΗ: μην είσαι συνοπτικός — γράψε ΕΚΤΕΝΕΣ περιεχόμενο με πολλές παραγράφους ανά ενότητα.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Τίτλος του εγγράφου"},
                    "sections_json": {"type": "string", "description": "JSON array από ενότητες — ΒΑΛΕ ΠΟΛΛΕΣ ενότητες (8-15) με ΕΚΤΕΝΕΣ περιεχόμενο. Κάθε section: {heading, content (αναλυτικό κείμενο πολλών προτάσεων), table:{headers,rows}, list:[], type:'text|table|list'}. Το content πρέπει να έχει πολλές παραγράφους, όχι μία πρόταση."},
                    "filename": {"type": "string", "description": "Όνομα αρχείου (π.χ. report.docx)"}
                },
                "required": ["title", "sections_json", "filename"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "generate_xlsx",
            "description": "Δημιούργησε ΕΠΑΓΓΕΛΜΑΤΙΚΟ Excel αρχείο (.xlsx) με φύλλα εργασίας, πίνακες, δεδομένα, headers, αριθμητικά στοιχεία. Ιδανικό για data reports, budgets, financials, προσφορές με αριθμούς. Αποθηκεύεται αυτόματα στο uploads.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Όνομα αρχείου (π.χ. data.xlsx)"},
                    "sheets_json": {"type": "string", "description": "JSON array φύλλων: [{name, headers:[], rows:[[]]}]"}
                },
                "required": ["filename", "sheets_json"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "generate_pptx",
            "description": "Δημιούργησε ΕΠΑΓΓΕΛΜΑΤΙΚΗ παρουσίαση PowerPoint (.pptx) με slides, τίτλους, bullet points και περιεχόμενο. Ιδανικό για pitches, παρουσιάσεις σε πελάτες, company profiles. Αποθηκεύεται αυτόματα στο uploads.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Τίτλος παρουσίασης"},
                    "slides_json": {"type": "string", "description": "JSON array slides: [{title, content, bullet_points:[]}]"},
                    "filename": {"type": "string", "description": "Όνομα αρχείου (π.χ. presentation.pptx)"}
                },
                "required": ["title", "slides_json", "filename"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "lookup_word",
            "description": "Αναζήτησε λέξη στα έγκυρα διαδικτυακά ελληνικά λεξικά (Λεξικό Τριανταφυλλίδη, greek-language.gr). ΧΡΗΣΙΜΟΠΟΙΗΣΕ το ΓΙΑ ΚΑΘΕ ΑΜΦΙΒΟΛΙΑ ορθογραφίας, γραμματικής, κλίσης, σημασίας, τονισμού, ετυμολογίας ή συντακτικού. Ιδανικό για να επιβεβαιώσεις τη σωστή γραφή λέξεων, ειδικά για λόγιο/επίσημο λεξιλόγιο, επιστημονικούς όρους, ξένες λέξεις, ομόηχα.",
            "parameters": {
                "type": "object",
                "properties": {
                    "word": {"type": "string", "description": "Η λέξη προς αναζήτηση (ελληνική ή ξένη)"}
                },
                "required": ["word"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_uploads",
            "description": "Λίστα αρχείων που έχεις ανεβάσει ή έχει λάβει ο agent. Διαβάζει το uploads directory σου για να δεις ποια αρχεία είναι διαθέσιμα.",
            "parameters": {
                "type": "object",
                "properties": {}
            }
        }
    },
]

def resolve_path(path, agent_id=None):
    p = os.path.realpath(os.path.expanduser(path))
    if not os.path.isabs(p):
        p = os.path.realpath(os.path.join(AION_DIR, p))
    allowed = [os.path.realpath(AION_DIR)]
    if not any(p.startswith(root + "/") or p == root for root in allowed):
        raise PermissionError(f"Access denied: {p} is outside AION directory")
    return p

import re as _re

def _extract_json_balanced(text, start_idx):
    """Extract a balanced JSON object starting from an opening brace at start_idx.
    Returns (json_str, end_idx) or (None, start_idx) on failure."""
    if start_idx >= len(text) or text[start_idx] != '{':
        return None, start_idx
    depth = 0
    in_str = False
    escape = False
    for i in range(start_idx, len(text)):
        ch = text[i]
        if escape:
            escape = False
            continue
        if ch == '\\' and in_str:
            escape = True
            continue
        if ch == '"' and not escape:
            in_str = not in_str
            continue
        if not in_str:
            if ch == '{':
                depth += 1
            elif ch == '}':
                depth -= 1
                if depth == 0:
                    return text[start_idx:i+1], i+1
    return None, start_idx

def parse_xml_tool_calls(text):
    """Parse XML/JSON tool calls from engine responses.
    Supports:
      - <|tool_call|> {"name": "...", "arguments": {...}}
      - <functioncall> {"name": "...", "arguments": {...}}
      - <function=name>...</function>
      - <invoke name='tool'>...</invoke>
    Returns (tool_calls_list, cleaned_text) or (None, text) if no XML found."""
    import re as _re2

    # Find all tool call markers and try to extract JSON after them
    markers = ['<|tool_call|>', '<functioncall>', '<invoke name=', '<function=']
    close_tags = ['</|tool_call|>', '</functioncall>', '</invoke>', '</function>']

    all_matches = []
    idx = 0
    while idx < len(text):
        best_marker = None
        best_pos = len(text)
        for marker, close_tag in zip(markers, close_tags):
            pos = text.find(marker, idx)
            if 0 <= pos < best_pos:
                best_pos = pos
                best_marker = (marker, close_tag)

        if best_marker is None or best_pos >= len(text):
            break

        marker, close_tag = best_marker
        start = best_pos + len(marker)

        # For <invoke name='xxx'> or <function=xxx>, skip to end of opening tag
        if marker in ('<invoke name=', '<function='):
            gt = text.find('>', start)
            if gt >= 0:
                start = gt + 1
            else:
                break

        # Skip whitespace
        while start < len(text) and text[start] in ' \t\n\r':
            start += 1

        # Try to extract balanced JSON
        if start < len(text) and text[start] == '{':
            json_str, end = _extract_json_balanced(text, start)
            if json_str:
                # Check for closing tag
                if close_tag:
                    close_idx = text.find(close_tag, start)
                    if close_idx >= 0:
                        end = close_idx + len(close_tag)
                    else:
                        # No closing tag found, use end of JSON
                        pass
                # For <invoke name='xxx'> format, extract name from opening tag
                full_name = ""
                if marker == '<invoke name=':
                    m2 = _re2.search(r'name=["\']([^"\']+)["\']', text[best_pos:best_pos+80])
                    if m2:
                        full_name = m2.group(1)

                all_matches.append((best_pos, end, text[start:end], full_name))
                idx = end
                continue

        # Fallback: find closing tag
        if close_tag:
            close_idx = text.find(close_tag, start)
            if close_idx >= 0:
                content = text[start:close_idx].strip()
                full_name = ""
                if marker == '<invoke name=':
                    m2 = _re2.search(r'name=["\']([^"\']+)["\']', text[best_pos:best_pos+80])
                    if m2:
                        full_name = m2.group(1)
                all_matches.append((best_pos, close_idx + len(close_tag), content, full_name))
                idx = close_idx + len(close_tag)
                continue

        idx = best_pos + len(marker)

    if not all_matches:
        return None, text

    # Sort by position (earliest first)
    all_matches.sort(key=lambda x: x[0])

    tool_calls = []
    for all_match in all_matches:
        if len(all_match) == 4:
            start_pos, end_pos, content, invoke_name = all_match
        else:
            start_pos, end_pos, content = all_match
            invoke_name = ""
        # Try to extract JSON
        json_start = content.find('{')
        if json_start >= 0:
            json_str, _ = _extract_json_balanced(content, json_start)
            if json_str:
                try:
                    data = json.loads(json_str)
                except:
                    data = None
                if data:
                    name = data.get("name", "") or data.get("function", "") or ""
                    args = data.get("arguments", {})
                    if isinstance(args, str):
                        try:
                            args = json.loads(args)
                        except:
                            args = {}
                    if name:
                        tool_calls.append({
                            "id": f"call_{name}_{len(tool_calls)}",
                            "type": "function",
                            "function": {"name": name, "arguments": json.dumps(args, ensure_ascii=False)}
                        })
                        continue

        # Fallback: parse key:value pairs (for old <invoke> format)
        name = invoke_name or ""
        args = {}
        for line in content.split('\n'):
            line = line.strip()
            if not name and ':' in line:
                name = line.split(':', 1)[0].strip()
            elif ':' in line:
                k, v = line.split(':', 1)
                args[k.strip()] = v.strip()
        if name:
            tool_calls.append({
                "id": f"call_{name}_{len(tool_calls)}",
                "type": "function",
                "function": {"name": name, "arguments": json.dumps(args, ensure_ascii=False)}
            })

    # Remove all matched patterns from text
    cleaned = text
    for m in reversed(sorted(all_matches, key=lambda x: x[0])):
        start_pos, end_pos = m[0], m[1]
        cleaned = cleaned[:start_pos] + cleaned[end_pos:]
    cleaned = cleaned.strip()

    return tool_calls, cleaned

def find_uploaded_file(fname):
    """Search for a file across all agent upload directories."""
    from agents import AGENTS
    upload_base = str(UPLOADS_DIR)
    for a in AGENTS + [{"id": "ceo"}]:
        candidate = os.path.join(upload_base, a["id"], fname)
        if os.path.exists(candidate):
            return candidate
    # Also try as-is
    if os.path.exists(fname):
        return fname
    return None

def execute_tool(name, args, agent_id="agent"):
    try:
        result = _execute_tool_impl(name, args, agent_id)
        log_activity(agent_id, name, args, result, True)
        return result
    except subprocess.TimeoutExpired:
        result = "Command timed out"
        log_activity(agent_id, name, args, result, False)
        return result
    except Exception as e:
        result = f"Error in {name}: {str(e)}"
        log_activity(agent_id, name, args, result, False)
        return result

def _execute_tool_impl(name, args, agent_id="agent"):
    try:
        if name == "read_file":
            p = resolve_path(args["path"], agent_id)
            if not os.path.exists(p):
                found = find_uploaded_file(os.path.basename(args["path"]))
                if found:
                    p = found
                else:
                    return f"File not found: {p}"
            if p.endswith(".docx"):
                try:
                    from docx import Document
                    doc = Document(p)
                    text = "\n".join(p.text for p in doc.paragraphs)
                    return text or "(κενό Word αρχείο)"
                except ImportError:
                    return "Το python-docx δεν είναι εγκατεστημένο. Τρέξε: pip3 install python-docx"
                except Exception as e:
                    return f"Σφάλμα ανάγνωσης Word: {e}"
            with open(p) as f:
                return f.read()
        elif name == "write_file":
            p = resolve_path(args["path"], agent_id)
            os.makedirs(os.path.dirname(p), exist_ok=True)
            with open(p, "w") as f:
                f.write(args["content"])
            # Auto-index into knowledge base
            try:
                from kb import index_file, _get_current_project
                project = _get_current_project()
                index_file(project, p, agent_id)
            except:
                pass
            return f"Written: {p}"
        elif name == "list_dir":
            p = resolve_path(args["path"], agent_id)
            if not os.path.isdir(p):
                return f"Not a directory: {p}"
            items = os.listdir(p)
            return "\n".join(sorted(items))
        elif name == "run_command":
            ALLOWED_AGENTS = ("ceo", "dev", "analytics")
            if agent_id not in ALLOWED_AGENTS:
                return f"❌ Ο agent '{agent_id}' δεν έχει δικαίωμα εκτέλεσης εντολών. Μόνο: {', '.join(ALLOWED_AGENTS)}"
            BLOCKED_PATTERNS = ["rm -rf /", "dd if=", "mkfs", "> /dev/", ":(){ :|:& };:", "chmod 777", "sudo ", "> /etc/"]
            command = args["command"]
            for pattern in BLOCKED_PATTERNS:
                if pattern in command:
                    return f"❌ Blocked: '{pattern}' δεν επιτρέπεται"
            timeout = args.get("timeout", 30)
            result = subprocess.run(
                command, shell=True, capture_output=True, text=True, timeout=timeout
            )
            out = result.stdout[-3000:] if len(result.stdout) > 3000 else result.stdout
            err = result.stderr[-1000:] if len(result.stderr) > 1000 else result.stderr
            return (out + ("\n---STDERR---\n" + err if err else "")) or "(no output)"
        elif name == "web_search":
            perplexity_key = os.environ.get("PERPLEXITY_API_KEY", "")
            if not perplexity_key:
                perplexity_key = "PERPLEXITY_KEY_REMOVED"
            resp = requests.post(
                "https://api.perplexity.ai/chat/completions",
                headers={"Authorization": f"Bearer {perplexity_key}", "Content-Type": "application/json"},
                json={
                    "model": "sonar",
                    "messages": [
                        {"role": "system", "content": "Search the web and provide accurate results."},
                        {"role": "user", "content": args["query"]}
                    ],
                    "max_tokens": 2000,
                },
                timeout=30,
            )
            if resp.status_code == 200:
                return resp.json()["choices"][0]["message"]["content"]
            return f"Search error: {resp.status_code} {resp.text[:200]}"
        elif name == "web_fetch":
            fmt = args.get("format", "markdown")
            last_err = ""
            for attempt in range(3):
                resp = requests.get(args["url"], timeout=30, headers={"User-Agent": "AIONCLAW/1.0"})
                if resp.status_code == 200:
                    return resp.text[:5000]
                if resp.status_code in (429, 503):
                    wait = 2 ** attempt
                    time.sleep(wait)
                    last_err = f"{resp.status_code} (retry {attempt+1}/3 after {wait}s)"
                else:
                    return f"Fetch error: {resp.status_code}"
            return f"Fetch error: {last_err}"
        elif name == "remember":
            mem = load_memory()
            facts = _ensure_facts(mem)
            facts[args["key"]] = {
                "value": args["value"],
                "agent": "user",
                "source": "user",
                "updated": datetime.now().isoformat(),
            }
            result = save_memory(mem)
            return f"Αποθηκεύτηκε: {args['key']} = {args['value']}" if result == "OK" else f"Σφάλμα: {result}"
        elif name == "recall":
            now = time.time()
            cache_key = args["key"]
            cached = _recall_cache.get(cache_key)
            if cached and (now - cached["ts"]) < RECALL_CACHE_TTL:
                return cached["result"]
            mem = load_memory()
            facts = _get_facts(mem)
            key = args["key"]
            exact = facts.get(key)
            if exact:
                val = exact["value"] if isinstance(exact, dict) else exact
                result = f"{key}: {val}"
                _recall_cache[cache_key] = {"result": result, "ts": now}
                return result
            matches = {}
            for k, v in facts.items():
                if key.lower() in k.lower():
                    val = v["value"] if isinstance(v, dict) else v
                    matches[k] = val
            if matches:
                result = "\n".join(f"{k}: {v}" for k, v in matches.items())
                _recall_cache[cache_key] = {"result": result, "ts": now}
                return result
            _recall_cache.pop(cache_key, None)
            return f"Δεν βρέθηκε: {key}"
        elif name == "list_memories":
            mem = load_memory()
            facts = _get_facts(mem)
            if not facts:
                return "Κενή μνήμη"
            lines = []
            for k, v in facts.items():
                val = v["value"] if isinstance(v, dict) else v
                lines.append(f"{k}: {val}")
            return "\n".join(lines)
        elif name == "read_leads":
            leads_file = str(LEADS_FILE)
            try:
                with open(leads_file) as f:
                    data = json.load(f)
                leads = data if isinstance(data, list) else data.get("leads", [])
                search = args.get("search", "").lower()
                limit = args.get("limit", 10)
                if search:
                    leads = [l for l in leads if search in json.dumps(l).lower()]
                leads = leads[:limit]
                if not leads:
                    return "Δεν βρέθηκαν leads"
                result = []
                for l in leads:
                    result.append(f"{l.get('name', l.get('company', '?'))} | {l.get('status', '?')} | {l.get('email', '')}")
                return "\n".join(result)
            except Exception as e:
                return f"Error reading leads: {e}"
        elif name == "save_lead":
            leads_dir = str(CRM_DIR / "leads")
            leads_file = os.path.join(leads_dir, "leads-database.json")
            os.makedirs(leads_dir, exist_ok=True)
            try:
                with open(leads_file) as f:
                    data = json.load(f)
            except:
                data = {"meta": {"schema": "AION CRM v1.0", "lastUpdated": "", "totalLeads": 0, "byStatus": {"incoming": 0, "qualified": 0, "contacted": 0, "converted": 0, "rejected": 0}}, "leads": []}
            if isinstance(data, list):
                data = {"meta": {}, "leads": data}
            new_lead = {
                "name": args["name"],
                "website": args.get("website", ""),
                "industry": args.get("industry", ""),
                "location": args.get("location", ""),
                "contact": args.get("contact", ""),
                "serviceNeeded": args.get("serviceNeeded", ""),
                "onlinePresence": args.get("onlinePresence", "unknown"),
                "source": args.get("source_url", "leadfinder"),
                "date": datetime.now().strftime("%Y-%m-%d"),
                "status": "incoming"
            }
            if args.get("notes"):
                new_lead["notes"] = args["notes"]
            data.setdefault("leads", []).append(new_lead)
            data["meta"]["lastUpdated"] = datetime.now().strftime("%Y-%m-%dT%H:%M")
            data["meta"]["totalLeads"] = len(data["leads"])
            existing = data["meta"].get("byStatus", {})
            existing.setdefault("incoming", 0)
            existing["incoming"] = sum(1 for l in data["leads"] if l.get("status") == "incoming")
            data["meta"]["byStatus"] = existing
            with open(leads_file, "w") as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
            return f"✓ Lead saved: {args['name']}"
        elif name == "get_time":
            return datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        elif name == "delegate_to_agent":
            from agents import AGENTS
            from collaboration import bus, run_sub_agent, save_to_agent_session
            from performance import get_eta
            agent_id = args["agent_id"]
            known_ids = {a["id"] for a in AGENTS}
            if agent_id not in known_ids:
                available = ", ".join(sorted(known_ids))
                return f"❌ Ο agent '{agent_id}' ΔΕΝ υπάρχει στο σύστημα. Διαθέσιμοι agents: {available}. Χρησιμοποίησε list_agents για να δεις την ομάδα."
            task = args["task"]
            context = args.get("context", "")
            estimated_seconds = get_eta(agent_id)
            started_at = datetime.now().isoformat()
            PERF_START[f"delegate_{agent_id}"] = time.time()
            bus.status(agent_id, True, "writing")
            bus.log("ceo", agent_id, "delegate", f"Ανάθεση: {task[:200]}...")
            # Progress: start
            bus.broadcast({
                "type": "task_progress",
                "agent_id": agent_id,
                "status": "started",
                "progress": 0,
                "message": f"🚀 {agent_id} ξεκινά... (~{estimated_seconds}s)",
                "estimated_seconds": estimated_seconds,
                "remaining_seconds": estimated_seconds,
                "started_at": started_at,
                "ts": datetime.now().isoformat(),
            })
            bus.broadcast({
                "type": "agent_thinking",
                "agent_id": agent_id,
                "status": "started",
                "thought": f"🚀 {agent_id}: ξεκινά εργασία (εκτίμ. {estimated_seconds}s)",
                "estimated_seconds": estimated_seconds,
                "started_at": started_at,
                "ts": datetime.now().isoformat(),
            })
            result = run_sub_agent(agent_id, task, context, format=args.get("format", "compact"))
            duration = time.time() - PERF_START.pop(f"delegate_{agent_id}", time.time())
            # Progress: done
            bus.broadcast({
                "type": "task_progress",
                "agent_id": agent_id,
                "status": "complete",
                "progress": 100,
                "message": f"✅ {agent_id} ολοκλήρωσε ({duration:.1f}s)",
                "estimated_seconds": estimated_seconds,
                "duration_s": round(duration, 1),
                "started_at": started_at,
                "ts": datetime.now().isoformat(),
            })
            bus.log(agent_id, "ceo", "result", result[:500])
            bus.status(agent_id, False, "has_response")
            store_collab_memory(agent_id, task, result)
            save_to_agent_session(agent_id, "default", f"Από CEO: {task}", result)
            # Broadcast real-time chat message to update agent's chat live
            bus.broadcast({
                "type": "agent_chat",
                "agent_id": agent_id,
                "session_id": "default",
                "exchange": [
                    {"role": "user", "content": f"Από CEO: {task}", "_aid": agent_id, "_sid": "default"},
                    {"role": "assistant", "content": result[:2000], "_aid": agent_id, "_sid": "default"}
                ]
            })
            return f"Αποτέλεσμα από {agent_id} ({duration:.1f}s):\n\n{result}"
        elif name == "parallel_delegate":
            from agents import AGENTS
            from collaboration import bus, run_sub_agent, save_to_agent_session
            from concurrent.futures import ThreadPoolExecutor, as_completed
            delegations = args["delegations"]
            synthesize = args.get("synthesize", False)
            known_ids = {a["id"] for a in AGENTS}
            validated = []
            for d in delegations:
                aid = d["agent_id"]
                if aid not in known_ids:
                    return f"❌ Ο agent '{aid}' ΔΕΝ υπάρχει στο σύστημα. Διαθέσιμοι: {', '.join(sorted(known_ids))}"
                validated.append((aid, d["task"], d.get("context", ""), d.get("format", "compact")))
            bus.log("ceo", ", ".join(a for a,_,_,_ in validated), "parallel_delegate", f"Παράλληλη ανάθεση σε {len(validated)} agents")
            total = len(validated)
            results = {}
            with ThreadPoolExecutor(max_workers=total) as pool:
                fut_map = {}
                for aid, task, ctx, fmt in validated:
                    bus.status(aid, True, "writing")
                    bus.broadcast({
                        "type": "agent_thinking",
                        "agent_id": aid,
                        "status": "started",
                        "thought": f"🚀 {aid}: ξεκινά παράλληλα με {total} agents"
                    })
                    fut = pool.submit(run_sub_agent, aid, task, ctx, format=fmt)
                    fut_map[fut] = (aid, task)
                for fut in as_completed(fut_map):
                    aid, task = fut_map[fut]
                    try:
                        result = fut.result(timeout=300)
                        results[aid] = result
                        bus.status(aid, False, "has_response")
                        bus.broadcast({
                            "type": "agent_thinking",
                            "agent_id": aid,
                            "status": "complete",
                            "thought": f"✅ {aid} ολοκλήρωσε παράλληλα"
                        })
                        store_collab_memory(aid, task, result)
                        save_to_agent_session(aid, "default", f"Παράλληλο από CEO: {task}", result)
                    except Exception as e:
                        results[aid] = f"❌ Σφάλμα: {e}"
                        bus.broadcast({
                            "type": "agent_thinking",
                            "agent_id": aid,
                            "status": "error",
                            "thought": f"❌ {aid} απέτυχε: {str(e)[:100]}"
                        })
            if synthesize:
                combined = "## Παράλληλα Αποτελέσματα\n\n"
                for aid, r in results.items():
                    combined += f"### {aid}\n{r[:2000]}\n\n"
                return combined
            parts = [f"**{aid}** ({len(r)} chars): {r[:500]}" for aid, r in results.items()]
            return f"Παράλληλη εκτέλεση {len(results)} agents:\n" + "\n---\n".join(parts)
        elif name == "send_to_agent":
            from agents import AGENTS
            from collaboration import bus, run_sub_agent, save_to_agent_session
            from performance import get_eta
            to_agent = args["agent_id"]
            known_ids = {a["id"] for a in AGENTS}
            if to_agent not in known_ids:
                available = ", ".join(sorted(known_ids))
                return f"❌ Ο agent '{to_agent}' ΔΕΝ υπάρχει στο σύστημα. Διαθέσιμοι agents: {available}. Χρησιμοποίησε list_agents για να δεις την ομάδα."
            msg = args["message"]
            ctx_extra = args.get("context", "")
            task_with_context = f"{msg}\n\nContext: {ctx_extra}" if ctx_extra else msg
            estimated_seconds = get_eta(to_agent)
            started_at = datetime.now().isoformat()
            PERF_START[f"send_{to_agent}"] = time.time()
            bus.status(to_agent, True, "writing")
            bus.log("ceo", to_agent, "forward", f"Μήνυμα: {msg[:200]}...")
            bus.broadcast({
                "type": "agent_thinking",
                "agent_id": to_agent,
                "status": "started",
                "thought": f"📨 {to_agent}: επεξεργάζεται μήνυμα (εκτίμ. {estimated_seconds}s)",
                "estimated_seconds": estimated_seconds,
                "started_at": started_at,
                "ts": datetime.now().isoformat(),
            })
            result = run_sub_agent(to_agent, task_with_context, format=args.get("format", "compact"))
            duration = time.time() - PERF_START.pop(f"send_{to_agent}", time.time())
            bus.status(to_agent, False, "has_response")
            save_to_agent_session(to_agent, "default", f"Από CEO: {msg}", result)
            bus.broadcast({
                "type": "agent_chat",
                "agent_id": to_agent,
                "session_id": "default",
                "exchange": [
                    {"role": "user", "content": f"Από CEO: {msg}", "_aid": to_agent, "_sid": "default"},
                    {"role": "assistant", "content": result[:2000], "_aid": to_agent, "_sid": "default"}
                ]
            })
            bus.log(to_agent, "ceo", "reply", result[:500])
            return f"Απάντηση από {to_agent} ({duration:.1f}s):\n\n{result}"
        elif name == "send_file_to_agent":
            from agents import AGENTS
            from collaboration import bus
            to_agent = args["agent_id"]
            known_ids = {a["id"] for a in AGENTS} | {"ceo"}
            if to_agent not in known_ids:
                available = ", ".join(sorted(known_ids))
                return f"❌ Ο agent '{to_agent}' ΔΕΝ υπάρχει. Διαθέσιμοι: {available}."
            src = resolve_path(args["file_path"])
            if not os.path.exists(src):
                found = find_uploaded_file(os.path.basename(args["file_path"]))
                if found:
                    src = found
                else:
                    return f"File not found: {args['file_path']}"
            dest_name = args.get("rename") or os.path.basename(src)
            dest_dir = os.path.join(str(UPLOADS_DIR), to_agent)
            os.makedirs(dest_dir, exist_ok=True)
            dest = os.path.join(dest_dir, dest_name)
            import shutil
            shutil.copy2(src, dest)
            fsize = os.path.getsize(dest)
            bus.log("agent", to_agent, "file_sent", f"Αρχείο: {dest_name} ({fsize} bytes)")
            bus.broadcast({
                "type": "agent_chat",
                "agent_id": to_agent,
                "session_id": "default",
                "exchange": [
                    {"role": "user", "content": f"📎 Λήψη αρχείου: {dest_name} ({fsize} bytes)", "_aid": to_agent, "_sid": "default"}
                ]
            })
            # Update agent_files in frontend via collab broadcast
            bus.broadcast({
                "type": "file_updated",
                "agent_id": to_agent,
                "filename": dest_name,
            })
            return f"Το αρχείο {dest_name} στάλθηκε στον {to_agent} ({fsize} bytes). Ο παραλήπτης μπορεί να το διαβάσει με read_file('{dest}')"
        elif name == "list_agents":
            from agents import AGENTS
            lines = [f"Σύστημα έχει {len(AGENTS)} agents:",
                     "─" * 40]
            for a in AGENTS:
                tools_list = ", ".join(a.get("tools", [])[:6])
                lines.append(f"  {a['icon']} {a['name']} ({a['id']})")
                lines.append(f"     Ρόλος: {a['role']}")
                lines.append(f"     Εργαλεία: {tools_list}")
            return "\n".join(lines)
        elif name == "query_kb":
            from kb import query_knowledge, format_kb_results
            q = args["query"]
            project = args.get("project", "")
            results = query_knowledge(project=project if project else None, query=q)
            return format_kb_results(results, q)
        elif name == "lookup_word":
            word = args["word"].strip()
            if not word:
                return "Δώσε μια λέξη προς αναζήτηση."
            try:
                safe_word = requests.utils.quote(word)
                # Try Triantafyllidis dictionary
                url = f"https://www.greek-language.gr/greekLang/modern_greek/tools/lexica/triantafyllides/search.html?lq={safe_word}"
                resp = requests.get(url, timeout=15, headers={"User-Agent": "AIONCLAW/1.0"})
                if resp.status_code == 200:
                    html = resp.text
                    # Find <div id="lemmas"> then extract <dl><dt> blocks
                    lemmas_match = _re.search(r'<div id="lemmas">(.*?)</div>\s*<form', html, _re.DOTALL)
                    if lemmas_match:
                        lemma_html = lemmas_match.group(1)
                        dt_blocks = _re.findall(r'<dt>(.*?)</dt>', lemma_html, _re.DOTALL)
                        texts = []
                        for dt in dt_blocks[:3]:
                            text = _re.sub(r'<[^>]+>', ' ', dt)
                            text = _re.sub(r'\s+', ' ', text).strip()
                            if text and len(text) > 10:
                                texts.append(text[:500])
                        if texts:
                            result = f"📖 Λεξικό Τριανταφυλλίδη — «{word}»:\n" + "\n\n".join(texts)
                            result += f"\n\n🔗 {url}"
                            return result
                # Try Academy of Athens dictionary as fallback
                aa_url = f"https://christikolexiko.academyofathens.gr/index.php/anazitisi?st={safe_word}"
                aa_resp = requests.get(aa_url, timeout=15, headers={"User-Agent": "AIONCLAW/1.0"})
                if aa_resp.status_code == 200:
                    aa_html = aa_resp.text
                    # Try various result containers
                    for pattern in [r'<div class="lexicon-result-item">(.*?)</div>',
                                    r'<div[^>]*class="[^"]*result[^"]*"[^>]*>(.*?)</div>']:
                        aa_blocks = _re.findall(pattern, aa_html, _re.DOTALL)
                        if aa_blocks:
                            texts = []
                            for b in aa_blocks[:3]:
                                text = _re.sub(r'<[^>]+>', ' ', b)
                                text = _re.sub(r'\s+', ' ', text).strip()
                                if text and len(text) > 10:
                                    texts.append(text[:500])
                            if texts:
                                result = f"📖 Χρηστικό Λεξικό Ακαδημίας Αθηνών — «{word}»:\n" + "\n\n".join(texts)
                                result += f"\n\n🔗 {aa_url}"
                                return result
                return f"Δεν βρέθηκε αποτέλεσμα για «{word}» στα λεξικά. Δοκίμασε διαφορετική γραφή.\n🔗 Λεξικό Τριανταφυλλίδη: {url}\n🔗 Ακαδημία Αθηνών: {aa_url}"
            except Exception as e:
                return f"Σφάλμα αναζήτησης για «{word}»: {e}\n🔗 Λεξικό Τριανταφυλλίδη: https://www.greek-language.gr/greekLang/modern_greek/tools/lexica/triantafyllides/search.html?lq={word.replace(' ', '%20')}"
        elif name == "list_uploads":
            upload_dir = os.path.join(str(UPLOADS_DIR), agent_id)
            if not os.path.isdir(upload_dir):
                return f"Δεν υπάρχουν ανεβασμένα αρχεία για τον agent '{agent_id}'."
            files = []
            for f in sorted(os.listdir(upload_dir)):
                fpath = os.path.join(upload_dir, f)
                size = os.path.getsize(fpath)
                mtime = datetime.fromtimestamp(os.path.getmtime(fpath)).strftime("%Y-%m-%d %H:%M")
                files.append(f"  {f} ({size:,} bytes, {mtime})")
            header = f"Αρχεία uploads για {agent_id} ({len(files)}):"
            return "\n".join([header] + files)
        elif name == "get_agent_history":
            target = args.get("agent", "")
            import glob
            session_dir = str(SESSIONS_DIR)
            # Search all project dirs for this agent's session files
            history_lines = []
            for root, _dirs, files in os.walk(session_dir):
                for fname in sorted(files):
                    if fname.startswith(f"{target}_") and fname.endswith(".json") and ":" not in fname:
                        fpath = os.path.join(root, fname)
                        try:
                            with open(fpath) as f:
                                data = json.load(f)
                            msgs = data.get("messages", [])
                            for m in msgs[-10:]:
                                role = m.get("role", "")
                                content = (m.get("content","") or "")[:200]
                                ts = m.get("ts","")[11:19] if m.get("ts") else ""
                                history_lines.append(f"[{ts}] {role}: {content}")
                        except:
                            pass
            if not history_lines:
                return f"Δεν βρέθηκε ιστορικό για τον agent '{target}'"
            return "Τελευταία μηνύματα:\n" + "\n".join(history_lines[-20:])
        elif name == "generate_docx":
            from collaboration import bus
            from docx import Document
            from docx.shared import Pt, Inches, Cm, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.table import WD_TABLE_ALIGNMENT
            doc = Document()
            style = doc.styles['Normal']
            style.font.name = 'Calibri'
            style.font.size = Pt(11)
            title = args.get("title", "Έγγραφο")
            p = doc.add_heading(title, level=0)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            sections = json.loads(args.get("sections_json", "[]"))
            for sec in sections:
                if sec.get("heading"):
                    doc.add_heading(sec["heading"], level=1)
                if sec.get("content"):
                    doc.add_paragraph(sec["content"])
                if sec.get("table"):
                    tbl = sec["table"]
                    rows = [tbl.get("headers", [])] + tbl.get("rows", [])
                    if rows:
                        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                        table.style = 'Light Grid Accent 1'
                        table.alignment = WD_TABLE_ALIGNMENT.CENTER
                        for ri, row_data in enumerate(rows):
                            for ci, val in enumerate(row_data):
                                cell = table.cell(ri, ci)
                                cell.text = str(val)
                                if ri == 0:
                                    for paragraph in cell.paragraphs:
                                        for run in paragraph.runs:
                                            run.bold = True
                if sec.get("list"):
                    for item in sec["list"]:
                        doc.add_paragraph(item, style='List Bullet')
            doc.add_paragraph("")
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            run = p.add_run(f"Generated by AIONCLAW — {datetime.now().strftime('%Y-%m-%d %H:%M')}")
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(128, 128, 128)
            fname = args.get("filename", "report.docx")
            out = os.path.join(str(UPLOADS_DIR), agent_id, fname)
            os.makedirs(os.path.dirname(out), exist_ok=True)
            doc.save(out)
            fsize = os.path.getsize(out)
            link = f"/api/files/download?path={AION_DIR}/aionclaw/uploads/{agent_id}/{fname}"
            from collaboration import bus
            bus.broadcast({"type": "file_updated", "agent_id": agent_id, "filename": fname})
            return f"✅ Word αρχείο: {fname} ({fsize:,} bytes)\n📎 Σύνδεσμος: {link}"
        elif name == "generate_xlsx":
            from collaboration import bus
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            sheets = json.loads(args.get("sheets_json", "[]"))
            for si, sh in enumerate(sheets):
                ws = wb.create_sheet(title=sh.get("name", f"Sheet{si+1}"))
                header_font = Font(bold=True, color="FFFFFF", size=11)
                header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                thin_border = Border(
                    left=Side(style='thin'), right=Side(style='thin'),
                    top=Side(style='thin'), bottom=Side(style='thin'))
                headers = sh.get("headers", [])
                if headers:
                    for ci, h in enumerate(headers, 1):
                        cell = ws.cell(row=1, column=ci, value=h)
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center')
                        cell.border = thin_border
                for ri, row in enumerate(sh.get("rows", []), 2 if headers else 1):
                    for ci, val in enumerate(row, 1):
                        cell = ws.cell(row=ri, column=ci, value=val)
                        cell.border = thin_border
                        cell.alignment = Alignment(horizontal='center')
                for ci in range(1, (len(headers) or 1) + 1):
                    ws.column_dimensions[get_column_letter(ci)].width = 18
            fname = args.get("filename", "data.xlsx")
            out = os.path.join(str(UPLOADS_DIR), agent_id, fname)
            os.makedirs(os.path.dirname(out), exist_ok=True)
            wb.save(out)
            fsize = os.path.getsize(out)
            link = f"/api/files/download?path={AION_DIR}/aionclaw/uploads/{agent_id}/{fname}"
            bus.broadcast({"type": "file_updated", "agent_id": agent_id, "filename": fname})
            return f"✅ Excel αρχείο: {fname} ({fsize:,} bytes)\n📎 Σύνδεσμος: {link}"
        elif name == "generate_pptx":
            from collaboration import bus
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor as PptRGB
            from pptx.enum.text import PP_ALIGN
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            title = args.get("title", "Παρουσίαση")
            slides = json.loads(args.get("slides_json", "[]"))
            for si, slide_data in enumerate(slides):
                if si == 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = PptRGB(44, 62, 80)
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = slide_data.get("title", title)
                    p.font.size = Pt(44)
                    p.font.color.rgb = PptRGB(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER
                    if slide_data.get("content"):
                        p2 = tf.add_paragraph()
                        p2.text = slide_data["content"]
                        p2.font.size = Pt(20)
                        p2.font.color.rgb = PptRGB(200, 200, 200)
                        p2.alignment = PP_ALIGN.CENTER
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title_shape = slide.shapes.title
                    title_shape.text = slide_data.get("title", f"Slide {si+1}")
                    content = slide_data.get("content", "")
                    bullets = slide_data.get("bullet_points", [])
                    if content or bullets:
                        body = slide.placeholders[1]
                        tf = body.text_frame
                        if content:
                            p = tf.paragraphs[0]
                            p.text = content
                            p.font.size = Pt(16)
                        for b in bullets:
                            p = tf.add_paragraph()
                            p.text = b
                            p.font.size = Pt(14)
                            p.level = 0
            fname = args.get("filename", "presentation.pptx")
            out = os.path.join(str(UPLOADS_DIR), agent_id, fname)
            os.makedirs(os.path.dirname(out), exist_ok=True)
            prs.save(out)
            fsize = os.path.getsize(out)
            link = f"/api/files/download?path={AION_DIR}/aionclaw/uploads/{agent_id}/{fname}"
            bus.broadcast({"type": "file_updated", "agent_id": agent_id, "filename": fname})
            return f"✅ PowerPoint αρχείο: {fname} ({fsize:,} bytes)\n📎 Σύνδεσμος: {link}"
        return f"Unknown tool: {name}"
    except subprocess.TimeoutExpired as e:
        raise e
    except Exception as e:
        raise e
